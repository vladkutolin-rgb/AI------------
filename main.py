"""
AI Presentation Generator
Ростелеком — Хакатон «Амурский Код» 2026
Полный продакшн-код с кэшированием, предпросмотром и отказоустойчивостью
"""

import json
import io
import uuid
import random
import base64
import traceback
import re
from typing import Optional, List, Dict, Any, Tuple
from fastapi import FastAPI, Form, UploadFile, File
from fastapi.responses import StreamingResponse, HTMLResponse, JSONResponse
import httpx
import pypdf
import docx2txt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import urllib3

# Отключаем предупреждения SSL для корпоративных сертификатов
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# ============================================================
# КОНФИГУРАЦИЯ ПРИЛОЖЕНИЯ
# ============================================================

app = FastAPI(
    title="AI Presentation Generator",
    description="Ростелеком — Интеллектуальный генератор презентаций",
    version="3.0.0"
)

# API-ключи и эндпоинты
API_TOKEN = "eyJhbGciOiJIUzM4NCJ9.eyJzY29wZXMiOlsibGxhbWEiLCJzZCIsInlhQXJ0Il0sInN1YiI6ImhhY2thdGhvbl8yNl8wNCIsImlhdCI6MTc3Njk0OTA3OCwiZXhwIjoxNzc3NjQwMjc4fQ.CbX06R1mPJmheZ8Tqc4HhaPdFDmYBKJgVhuxdJG5wPUfQArFZAroHrTjz2Q5SGuF"
BASE_URL = "https://ai.rt.ru/api/1.0"

# Эндпоинты API Ростелекома
LLM_ENDPOINT = f"{BASE_URL}/llama/chat"
YA_IMAGE_ENDPOINT = f"{BASE_URL}/ya/image"
SD_IMAGE_ENDPOINT = f"{BASE_URL}/sd/img"
DOWNLOAD_ENDPOINT = f"{BASE_URL}/download"

# Параметры слайдов
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5
DEFAULT_SLIDES_COUNT = 6
MAX_SLIDES = 20
MIN_SLIDES = 1

# ============================================================
# ХРАНИЛИЩЕ СЛАЙДОВ (кэш для соответствия предпросмотра и PPTX)
# ============================================================
slides_cache: Dict[str, List[Dict[str, Any]]] = {}

# ============================================================
# ДИЗАЙН-БИБЛИОТЕКИ
# ============================================================

# Цветовые палитры по настроениям
COLOR_PALETTES: Dict[str, Dict[str, Any]] = {
    "tech": {
        "name": "Технологичный",
        "gradients": [
            ("#1A1A2E", "#16213E"),
            ("#0F3460", "#1A1A2E"),
            ("#16213E", "#0F3460"),
            ("#1A1A2E", "#0F3460"),
            ("#16213E", "#1A1A2E"),
        ],
        "accents": ["#E94560", "#00D2FF", "#7B2FBE", "#FF6B6B", "#4ECDC4"],
        "is_dark": True,
    },
    "nature": {
        "name": "Природный",
        "gradients": [
            ("#F0F7F4", "#E8F5E9"),
            ("#E8F5E9", "#F1F8E9"),
            ("#F5FBF0", "#E8F0E8"),
            ("#F1F8E9", "#F0F7F4"),
            ("#FAFFF8", "#F0F7F0"),
        ],
        "accents": ["#2D6A4F", "#40916C", "#52B788", "#95D5B2", "#1B4332"],
        "is_dark": False,
    },
    "warm": {
        "name": "Тёплый",
        "gradients": [
            ("#FFF5F0", "#FFF0E6"),
            ("#FFF8F5", "#FFF3ED"),
            ("#FFFBF7", "#FFF5F0"),
            ("#FFF3ED", "#FFF8F5"),
            ("#FFF0E6", "#FFF5F0"),
        ],
        "accents": ["#FF6B6B", "#FF8E72", "#FFAA85", "#E8523A", "#FFB347"],
        "is_dark": False,
    },
    "cool": {
        "name": "Холодный",
        "gradients": [
            ("#F0F4FF", "#E8EEFF"),
            ("#F5F8FF", "#EDF2FF"),
            ("#FAFBFF", "#F0F4FF"),
            ("#EDF2FF", "#F5F8FF"),
            ("#E8EEFF", "#F0F4FF"),
        ],
        "accents": ["#1B3A5C", "#2D6B9F", "#69A2D6", "#A8D0F0", "#4A90D9"],
        "is_dark": False,
    },
    "luxury": {
        "name": "Премиум",
        "gradients": [
            ("#1A1A1A", "#2D2D2D"),
            ("#0D0D0D", "#1A1A1A"),
            ("#2D2D2D", "#1A1A1A"),
            ("#1A1A1A", "#0D0D0D"),
            ("#2D2D2D", "#0D0D0D"),
        ],
        "accents": ["#C6A962", "#DFD0A5", "#F5F0E8", "#D4AF37", "#B8860B"],
        "is_dark": True,
    },
    "vivid": {
        "name": "Яркий",
        "gradients": [
            ("#FFFFFF", "#FFF5F5"),
            ("#FFF8F8", "#FFF0F0"),
            ("#FFFBFB", "#FFF5F5"),
            ("#FFF5F5", "#FFFFFF"),
            ("#FFF0F0", "#FFF8F8"),
        ],
        "accents": ["#FF006E", "#FB5607", "#FFBE0B", "#8338EC", "#3A86FF"],
        "is_dark": False,
    },
    "minimal": {
        "name": "Минималистичный",
        "gradients": [
            ("#FAFAFA", "#FFFFFF"),
            ("#F5F5F7", "#FAFAFA"),
            ("#FFFFFF", "#F8F8FA"),
            ("#FAFAFA", "#F5F5F7"),
            ("#F8F8FA", "#FFFFFF"),
        ],
        "accents": ["#333333", "#555555", "#777777", "#999999", "#BBBBBB"],
        "is_dark": False,
    },
    "rostelcom": {
        "name": "Ростелеком",
        "gradients": [
            ("#FFFFFF", "#F8F5FF"),
            ("#F8F5FF", "#FFFFFF"),
            ("#FFFFFF", "#F5F0FF"),
            ("#FAF8FF", "#FFFFFF"),
            ("#F8F5FF", "#FAF8FF"),
        ],
        "accents": ["#5C2D91", "#FF6B00", "#7B4FBF", "#E85D04", "#9B5DE5"],
        "is_dark": False,
    },
}

# Стили шрифтов
FONT_LIBRARY: Dict[str, Dict[str, Any]] = {
    "modern": {"name": "Calibri", "title_size": 34, "title_weight": True},
    "classic": {"name": "Georgia", "title_size": 32, "title_weight": True},
    "tech": {"name": "Consolas", "title_size": 30, "title_weight": True},
    "minimal": {"name": "Arial", "title_size": 36, "title_weight": False},
}

# ============================================================
# HTTP КЛИЕНТЫ
# ============================================================

def create_sync_client() -> httpx.Client:
    """Создаёт синхронный HTTP клиент"""
    return httpx.Client(verify=False, timeout=180)

def create_async_client() -> httpx.AsyncClient:
    """Создаёт асинхронный HTTP клиент"""
    return httpx.AsyncClient(verify=False, timeout=60)

# ============================================================
# ОБРАБОТЧИКИ ДОКУМЕНТОВ
# ============================================================

def extract_pdf_text(file_bytes: bytes) -> str:
    """Извлекает текст из PDF с обработкой ошибок"""
    try:
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                parts.append(text.strip())
        result = " ".join(parts)
        return result[:5000]  # Ограничиваем длину
    except Exception as e:
        print(f"[PDF ERROR] {e}")
        return ""

def extract_docx_text(file_bytes: bytes) -> str:
    """Извлекает текст из DOCX с обработкой ошибок"""
    try:
        text = docx2txt.process(io.BytesIO(file_bytes))
        return (text or "")[:5000]  # Ограничиваем длину
    except Exception as e:
        print(f"[DOCX ERROR] {e}")
        return ""

# ============================================================
# ЦВЕТОВЫЕ УТИЛИТЫ
# ============================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Конвертирует HEX в RGBColor"""
    if not hex_color or not isinstance(hex_color, str):
        return RGBColor(0x99, 0x99, 0x99)
    
    hex_color = hex_color.lstrip("#").strip()
    
    try:
        if len(hex_color) == 6:
            return RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
        elif len(hex_color) == 3:
            return RGBColor(
                int(hex_color[0] * 2, 16),
                int(hex_color[1] * 2, 16),
                int(hex_color[2] * 2, 16)
            )
    except Exception:
        pass
    
    return RGBColor(0x99, 0x99, 0x99)

def is_dark_background(hex_color: str) -> bool:
    """Определяет, тёмный ли цвет фона"""
    try:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        # Формула воспринимаемой яркости
        brightness = (r * 299 + g * 587 + b * 114) / 1000
        return brightness < 128
    except Exception:
        return False

def calculate_text_color(bg_hex: str) -> Tuple[RGBColor, RGBColor]:
    """Вычисляет контрастные цвета текста для фона"""
    if is_dark_background(bg_hex):
        title_color = RGBColor(0xFF, 0xFF, 0xFF)
        body_color = RGBColor(0xDD, 0xDD, 0xEE)
    else:
        title_color = RGBColor(0x1A, 0x1A, 0x2E)
        body_color = RGBColor(0x3A, 0x3A, 0x4A)
    return title_color, body_color

def calculate_font_size(text: str, max_size: int = 24, min_size: int = 12) -> int:
    """Автоматически подбирает размер шрифта под длину текста"""
    length = len(text)
    if length < 80:
        return max_size
    elif length < 150:
        return max(20, min_size)
    elif length < 250:
        return max(18, min_size)
    elif length < 400:
        return max(15, min_size)
    else:
        return min_size

def image_to_base64_url(img_bytes: Optional[bytes]) -> Optional[str]:
    """Конвертирует байты изображения в data URL"""
    if not img_bytes:
        return None
    return f"data:image/png;base64,{base64.b64encode(img_bytes).decode()}"

# ============================================================
# ИНТЕЛЛЕКТУАЛЬНЫЙ АНАЛИЗ ЗАПРОСА
# ============================================================

def analyze_prompt(prompt: str) -> Dict[str, Any]:
    """
    Анализирует запрос пользователя и извлекает:
    - настроение (mood)
    - количество слайдов
    - стиль шрифта
    - особые пожелания
    """
    prompt_lower = prompt.lower()
    
    # Определяем настроение
    mood = "minimal"
    if any(w in prompt_lower for w in ["ростелеком", "rt", "рт", "связь", "провайдер"]):
        mood = "rostelcom"
    elif any(w in prompt_lower for w in ["тёмн", "темн", "dark", "tech", "техн", "кибер", "хакер"]):
        mood = "tech"
    elif any(w in prompt_lower for w in ["природ", "зелён", "зелен", "green", "nature", "эко", "лес"]):
        mood = "nature"
    elif any(w in prompt_lower for w in ["тёпл", "тепл", "warm", "уют", "оранж", "солнц"]):
        mood = "warm"
    elif any(w in prompt_lower for w in ["холод", "син", "blue", "cool", "спокой", "лёд"]):
        mood = "cool"
    elif any(w in prompt_lower for w in ["luxury", "золот", "чёрн", "черн", "премиум", "дорог", "элит"]):
        mood = "luxury"
    elif any(w in prompt_lower for w in ["ярк", "vivid", "цвет", "радуг", "неон"]):
        mood = "vivid"
    
    # Извлекаем количество слайдов
    slides_count = DEFAULT_SLIDES_COUNT
    match = re.search(r'(\d+)\s*(?:слайд|slide|slajd)', prompt_lower)
    if match:
        count = int(match.group(1))
        slides_count = max(MIN_SLIDES, min(MAX_SLIDES, count))
    
    # Определяем стиль шрифта
    font_style = "modern"
    if any(w in prompt_lower for w in ["classic", "класси", "georgia", "сериф", "строг"]):
        font_style = "classic"
    elif any(w in prompt_lower for w in ["tech", "console", "моно", "код", "програм"]):
        font_style = "tech"
    elif any(w in prompt_lower for w in ["minimal", "минимал", "прост", "чист"]):
        font_style = "minimal"
    
    return {
        "mood": mood,
        "slides_count": slides_count,
        "font_style": font_style,
        "palette": COLOR_PALETTES.get(mood, COLOR_PALETTES["minimal"]),
        "font": FONT_LIBRARY.get(font_style, FONT_LIBRARY["modern"]),
    }

# ============================================================
# LLM — ГЕНЕРАЦИЯ СТРУКТУРЫ ПРЕЗЕНТАЦИИ
# ============================================================

def build_llm_prompt(document_text: str, user_prompt: str, analysis: Dict[str, Any]) -> str:
    """Строит системный промпт для LLM"""
    palette = analysis["palette"]
    mood_name = palette["name"]
    slides_count = analysis["slides_count"]
    
    return f"""Ты — арт-директор презентаций. Создай структуру из {slides_count} слайдов.

СТИЛЬ: {mood_name}
ФОН: {"тёмный" if palette['is_dark'] else "светлый"}

ДЛЯ КАЖДОГО СЛАЙДА ОПРЕДЕЛИ:
- title: заголовок (до 6 слов, яркий и цепляющий)
- content: основной текст (25-40 слов, СОДЕРЖАТЕЛЬНО, с фактами и конкретикой)
- image_keywords: 3-4 английских слова для поиска релевантной картинки
- bg_color1, bg_color2: два HEX-цвета для градиентного фона
- accent_color: один HEX-цвет для акцентов

ВАЖНО: Используй цвета из {mood_name} палитры.

ОТВЕТЬ ТОЛЬКО JSON, БЕЗ КОММЕНТАРИЕВ:
{{"slides":[{{"title":"Заголовок","content":"Текст 25-40 слов","image_keywords":"english keywords","bg_color1":"#HEX","bg_color2":"#HEX","accent_color":"#HEX"}}]}}

ЗАПРОС: {user_prompt}
ДОКУМЕНТ: {document_text[:2500] if document_text else "нет"}"""

def call_llm_api(document_text: str, user_prompt: str) -> Dict[str, Any]:
    """Отправляет запрос к LLM и получает структуру презентации"""
    analysis = analyze_prompt(user_prompt)
    headers = {
        "Authorization": f"Bearer {API_TOKEN}",
        "Content-Type": "application/json"
    }
    
    system_msg = build_llm_prompt(document_text, user_prompt, analysis)
    
    payload = {
        "uuid": str(uuid.uuid4()),
        "chat": {
            "model": "Qwen/Qwen2.5-72B-Instruct",
            "user_message": system_msg,
            "contents": [{"type": "text", "text": system_msg}],
            "system_prompt": "Ты — профессиональный дизайнер презентаций. Отвечай ТОЛЬКО валидным JSON. Никаких пояснений.",
            "max_new_tokens": 5000,
            "temperature": 0.9,
        }
    }
    
    try:
        with create_sync_client() as client:
            response = client.post(LLM_ENDPOINT, json=payload, headers=headers)
        
        if response.status_code != 200:
            raise Exception(f"API error {response.status_code}: {response.text[:200]}")
        
        data = response.json()
        raw_content = data[0]["message"]["content"]
        
        # Очистка от markdown
        if "```json" in raw_content:
            raw_content = raw_content.split("```json")[1].split("```")[0]
        elif "```" in raw_content:
            raw_content = raw_content.split("```")[1].split("```")[0]
        
        # Извлечение JSON
        if "{" in raw_content and "}" in raw_content:
            start = raw_content.index("{")
            end = raw_content.rindex("}") + 1
            raw_content = raw_content[start:end]
        
        result = json.loads(raw_content)
        
        if "slides" not in result or not result["slides"]:
            raise ValueError("Empty slides array")
        
        return result
    
    except Exception as e:
        print(f"[LLM ERROR] {e}")
        traceback.print_exc()
        return generate_fallback_slides(user_prompt, analysis)

def generate_fallback_slides(user_prompt: str, analysis: Dict[str, Any]) -> Dict[str, Any]:
    """Генерирует запасные слайды если LLM недоступна"""
    palette = analysis["palette"]
    count = analysis["slides_count"]
    slides = []
    
    topics = [
        "Обзор темы", "Ключевые факты", "Статистика и цифры",
        "Преимущества", "Примеры", "Тренды", "Прогнозы",
        "Рекомендации", "Выводы", "Следующие шаги"
    ]
    
    for i in range(count):
        gradient = random.choice(palette["gradients"])
        accent = random.choice(palette["accents"])
        topic = topics[i % len(topics)]
        
        slides.append({
            "title": f"{topic}",
            "content": f"Ключевая информация по теме «{user_prompt[:60]}...». {topic.lower()} — важный аспект для понимания общей картины.",
            "image_keywords": f"abstract {palette['name'].lower()} modern",
            "bg_color1": gradient[0],
            "bg_color2": gradient[1],
            "accent_color": accent,
        })
    
    return {"slides": slides}

# ============================================================
# ГЕНЕРАЦИЯ ИЗОБРАЖЕНИЙ
# ============================================================

async def generate_image_ya_art(client: httpx.AsyncClient, prompt: str, headers: Dict) -> Tuple[Optional[str], str]:
    """Генерация через YandexArt"""
    payload = {
        "uuid": str(uuid.uuid4()),
        "image": {
            "request": prompt,
            "seed": random.randint(1, 99999),
            "translate": False,
            "model": "yandex-art",
            "aspect": "16:9"
        }
    }
    try:
        resp = await client.post(YA_IMAGE_ENDPOINT, json=payload, headers=headers)
        if resp.status_code == 200:
            data = resp.json()
            if data and len(data) > 0 and "message" in data[0]:
                return data[0]["message"]["id"], "yaArt"
    except Exception as e:
        print(f"[YA ART ERROR] {e}")
    return None, "yaArt"

async def generate_image_sd(client: httpx.AsyncClient, prompt: str, headers: Dict) -> Tuple[Optional[str], str]:
    """Генерация через Stable Diffusion"""
    payload = {
        "uuid": str(uuid.uuid4()),
        "sdImage": {
            "request": prompt,
            "seed": random.randint(1, 99999),
            "translate": False
        }
    }
    try:
        resp = await client.post(SD_IMAGE_ENDPOINT, json=payload, headers=headers)
        if resp.status_code == 200:
            data = resp.json()
            if data and len(data) > 0 and "message" in data[0]:
                return data[0]["message"]["id"], "sd"
    except Exception as e:
        print(f"[SD ERROR] {e}")
    return None, "sd"

async def download_generated_image(client: httpx.AsyncClient, image_id: str, 
                                   service_type: str, headers: Dict) -> Optional[bytes]:
    """Скачивает сгенерированное изображение"""
    try:
        url = f"{DOWNLOAD_ENDPOINT}?id={image_id}&serviceType={service_type}&imageType=png"
        resp = await client.get(url, headers=headers)
        if resp.status_code == 200:
            return resp.content
    except Exception as e:
        print(f"[DOWNLOAD ERROR] {e}")
    return None

async def generate_slide_image(keywords: str) -> Optional[bytes]:
    """Генерирует изображение для слайда (пробует YaArt, затем SD)"""
    if not keywords or not keywords.strip():
        return None
    
    headers = {
        "Authorization": f"Bearer {API_TOKEN}",
        "Content-Type": "application/json"
    }
    
    image_prompt = f"professional {keywords}, high quality, modern design, artistic, 4k"
    
    try:
        async with create_async_client() as client:
            # Пробуем YandexArt
            image_id, service = await generate_image_ya_art(client, image_prompt, headers)
            
            # Если не получилось — Stable Diffusion
            if not image_id:
                image_id, service = await generate_image_sd(client, image_prompt, headers)
            
            # Скачиваем
            if image_id:
                return await download_generated_image(client, image_id, service, headers)
    except Exception as e:
        print(f"[IMAGE GEN ERROR] {e}")
    
    return None

# ============================================================
# СБОРКА PPTX
# ============================================================

def add_gradient_background_to_slide(slide, color1_hex: str, color2_hex: str) -> None:
    """Добавляет градиентный фон на слайд"""
    try:
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_angle = random.choice([135, 180, 225, 270])
        stops = fill.gradient_stops
        stops[0].color.rgb = hex_to_rgb(color1_hex)
        stops[0].position = 0.0
        stops[1].color.rgb = hex_to_rgb(color2_hex)
        stops[1].position = 1.0
    except Exception:
        # Fallback: сплошной цвет
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color1_hex)

def add_decorative_elements(slide, accent_rgb: RGBColor, index: int) -> None:
    """Добавляет декоративные фигуры на слайд"""
    # Маленький круг в углу
    circle_size = random.uniform(0.2, 0.5)
    circle_x = random.choice([-0.1, 12.8, random.uniform(11.5, 12.5)])
    circle_y = random.choice([-0.1, 6.8, random.uniform(5.5, 7.0)])
    
    shape = slide.shapes.add_shape(
        9,  # Овал
        Inches(circle_x), Inches(circle_y),
        Inches(circle_size), Inches(circle_size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_rgb
    shape.line.fill.background()

def build_single_pptx_slide(slide, slide_data: Dict[str, Any], index: int) -> None:
    """Собирает один слайд презентации"""
    
    # Извлекаем данные
    title_text = slide_data.get("title", f"Слайд {index + 1}")
    content_text = slide_data.get("content", "")
    bg_color1 = slide_data.get("bg_color1", "#F5F5F7")
    bg_color2 = slide_data.get("bg_color2", "#FFFFFF")
    accent_hex = slide_data.get("accent_color", "#5C2D91")
    image_bytes = slide_data.get("image_bytes")
    
    # Конвертируем цвета
    accent_rgb = hex_to_rgb(accent_hex)
    title_rgb, body_rgb = calculate_text_color(bg_color1)
    
    # Фон
    add_gradient_background_to_slide(slide, bg_color1, bg_color2)
    
    # Декоративные элементы
    add_decorative_elements(slide, accent_rgb, index)
    
    # Номер слайда
    num_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(0.6), Inches(0.35))
    num_para = num_box.text_frame.paragraphs[0]
    num_para.text = f"0{index + 1}" if index < 9 else str(index + 1)
    num_para.font.size = Pt(11)
    num_para.font.color.rgb = accent_rgb
    num_para.font.bold = True
    
    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(8.0), Inches(1.1))
    title_box.text_frame.word_wrap = True
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = title_rgb
    
    # Линия под заголовком
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.8), Inches(2.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_rgb
    line.line.fill.background()
    
    # Основной текст
    font_size = calculate_font_size(content_text)
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(7.5), Inches(4.5))
    content_box.text_frame.word_wrap = True
    content_para = content_box.text_frame.paragraphs[0]
    content_para.text = content_text
    content_para.font.size = Pt(font_size)
    content_para.font.color.rgb = body_rgb
    content_para.line_spacing = Pt(font_size * 1.6)
    
    # Картинка справа
    if image_bytes:
        try:
            img_stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(
                img_stream,
                Inches(8.6), Inches(1.5),
                Inches(4.2), Inches(5.0)
            )
        except Exception as e:
            print(f"[IMAGE INSERT ERROR] {e}")

def build_complete_pptx(slides_data: List[Dict[str, Any]]) -> io.BytesIO:
    """Собирает полную PPTX презентацию из списка слайдов"""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой макет
        build_single_pptx_slide(slide, slide_data, i)
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ============================================================
# API ЭНДПОИНТЫ
# ============================================================

@app.post("/generate")
async def generate_presentation(
    prompt: str = Form(...),
    file: UploadFile = File(None),
    download: str = Form("false"),
    cache_id: str = Form("")
):
    """
    Главный эндпоинт генерации презентаций.
    
    Параметры:
    - prompt: текстовое описание презентации
    - file: опциональный PDF/DOCX файл
    - download: если "true" — возвращает PPTX файл
    - cache_id: ID кэшированных слайдов (чтобы не перегенерировать)
    """
    
    # Если запрос на скачивание с кэшем — отдаём готовый PPTX
    if download == "true" and cache_id and cache_id in slides_cache:
        cached_slides = slides_cache[cache_id]
        pptx_bytes = build_complete_pptx(cached_slides)
        return StreamingResponse(
            pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=presentation.pptx"}
        )
    
    # Извлечение текста из загруженного файла
    document_text = ""
    if file and file.filename:
        file_content = await file.read()
        filename_lower = file.filename.lower()
        if filename_lower.endswith(".pdf"):
            document_text = extract_pdf_text(file_content)
        elif filename_lower.endswith(".docx"):
            document_text = extract_docx_text(file_content)
    
    # Генерация структуры через LLM
    llm_result = call_llm_api(document_text, prompt)
    slides = llm_result.get("slides", [])
    
    if not slides:
        # Экстренный fallback
        analysis = analyze_prompt(prompt)
        slides = generate_fallback_slides(prompt, analysis)["slides"]
    
    # Генерация изображений для каждого слайда
    for slide in slides:
        keywords = slide.get("image_keywords", "abstract modern")
        img_bytes = await generate_slide_image(keywords)
        slide["image_bytes"] = img_bytes
        slide["image_url"] = image_to_base64_url(img_bytes)
    
    # Сохраняем в кэш
    new_cache_id = str(uuid.uuid4())
    slides_cache[new_cache_id] = slides
    
    # Если запрошено скачивание
    if download == "true":
        pptx_bytes = build_complete_pptx(slides)
        return StreamingResponse(
            pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=presentation.pptx"}
        )
    
    # Формируем данные для предпросмотра
    preview_slides = []
    for s in slides:
        preview_slides.append({
            "title": s.get("title", ""),
            "content": s.get("content", ""),
            "bg_color1": s.get("bg_color1", "#F5F5F7"),
            "bg_color2": s.get("bg_color2", "#FFFFFF"),
            "accent_color": s.get("accent_color", "#5C2D91"),
            "image_url": s.get("image_url"),
        })
    
    return JSONResponse({
        "slides": preview_slides,
        "cache_id": new_cache_id,
        "count": len(preview_slides)
    })


@app.get("/")
async def serve_html():
    """Отдаёт HTML-страницу"""
    return HTMLResponse(open("index.html", encoding="utf-8").read())


@app.get("/health")
async def health_check():
    """Проверка работоспособности сервиса"""
    return {
        "status": "ok",
        "service": "AI Presentation Generator",
        "version": "3.0.0",
        "cached_presentations": len(slides_cache)
    }